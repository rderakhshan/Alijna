#!/usr/bin/env python
# coding=utf-8
# Copyright 2025 The OPPO Inc. Personal AI team. All rights reserved.
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

import base64
import copy
import html
import json
import mimetypes
import os
import re
import shutil
import subprocess
import sys
import tempfile
import traceback
import zipfile
from typing import Any, Dict, List, Optional, Union
from urllib.parse import parse_qs, quote, unquote, urlparse, urlunparse

import mammoth
import markdownify
import pandas as pd
import pdfminer.high_level
import pptx
import puremagic
import pydub
import requests
import speech_recognition as sr
from bs4 import BeautifulSoup
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api.formatters import SRTFormatter


class _CustomMarkdownify(markdownify.MarkdownConverter):
    """Custom Markdown converter with enhanced HTML-to-Markdown conversion."""
    
    def __init__(self, **options: Any):
        options["heading_style"] = options.get("heading_style", markdownify.ATX)
        super().__init__(**options)

    def convert_hn(self, n: int, el: Any, text: str, convert_as_inline: bool) -> str:
        if not convert_as_inline and not re.search(r"^\n", text):
            return "\n" + super().convert_hn(n, el, text, convert_as_inline)
        return super().convert_hn(n, el, text, convert_as_inline)

    def convert_a(self, el: Any, text: str, convert_as_inline: bool):
        prefix, suffix, text = markdownify.chomp(text)
        if not text:
            return ""
            
        href = el.get("href")
        title = el.get("title")

        if href:
            try:
                parsed_url = urlparse(href)
                if parsed_url.scheme and parsed_url.scheme.lower() not in ["http", "https", "file"]:
                    return f"{prefix}{text}{suffix}"
                href = urlunparse(parsed_url._replace(path=quote(unquote(parsed_url.path))))
            except ValueError:
                return f"{prefix}{text}{suffix}"

        if (self.options["autolinks"] and text.replace(r"\_", "_") == href 
            and not title and not self.options["default_title"]):
            return f"<{href}>"
            
        if self.options["default_title"] and not title:
            title = href
            
        escaped_title = title.replace('"', r'\"') if title else ""
        title_part = f' "{escaped_title}"' if title else ""

        return f"{prefix}[{text}]({href}{title_part}){suffix}" if href else text

    def convert_img(self, el: Any, text: str, convert_as_inline: bool) -> str:
        alt = el.attrs.get("alt", "") or ""
        src = el.attrs.get("src", "") or ""
        title = el.attrs.get("title", "") or ""
        
        if convert_as_inline and el.parent.name not in self.options["keep_inline_images_in"]:
            return alt

        if src.startswith("data:"):
            src = src.split(",")[0] + "..."
            
        escaped_title = title.replace('"', r'\"') if title else ""
        title_part = f' "{escaped_title}"' if title else ""
        
        return f"![{alt}]({src}{title_part})"

    def convert_soup(self, soup: Any) -> str:
        return super().convert_soup(soup)


class DocumentConverterResult:
    """Result of document conversion containing title and text content."""
    
    def __init__(self, title: Optional[str] = None, text_content: str = ""):
        self.title = title
        self.text_content = text_content


class DocumentConverter:
    """Abstract base class for document converters."""
    
    def convert(self, local_path: str, **kwargs: Any) -> Optional[DocumentConverterResult]:
        raise NotImplementedError()


class PlainTextConverter(DocumentConverter):
    """Converter for plain text files."""
    
    def convert(self, local_path: str, **kwargs: Any) -> Optional[DocumentConverterResult]:
        content_type, _ = mimetypes.guess_type("__placeholder" + kwargs.get("file_extension", ""))

        text_content = ""
        with open(local_path, "rt", encoding="utf-8") as fh:
            text_content = fh.read()
            
        return DocumentConverterResult(title=None, text_content=text_content)


class HtmlConverter(DocumentConverter):
    """Converter for HTML files."""
    
    def convert(self, local_path: str, **kwargs: Any) -> Optional[DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None

        with open(local_path, "rt", encoding="utf-8") as fh:
            return self._convert(fh.read())

    def _convert(self, html_content: str) -> Optional[DocumentConverterResult]:
        soup = BeautifulSoup(html_content, "html.parser")

        for script in soup(["script", "style"]):
            script.extract()

        body_elm = soup.find("body")
        if body_elm:
            webpage_text = _CustomMarkdownify().convert_soup(body_elm)
        else:
            webpage_text = _CustomMarkdownify().convert_soup(soup)

        title = None if soup.title is None else soup.title.string
        return DocumentConverterResult(title=title, text_content=str(webpage_text))


class WikipediaConverter(DocumentConverter):
    """Specialized converter for Wikipedia pages."""
    
    def convert(self, local_path: str, **kwargs: Any) -> Optional[DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None
            
        url = kwargs.get("url", "")
        if not re.search(r"^https?:\/\/[a-zA-Z]{2,3}\.wikipedia.org\/", url):
            return None

        with open(local_path, "rt", encoding="utf-8") as fh:
            soup = BeautifulSoup(fh.read(), "html.parser")

        for script in soup(["script", "style"]):
            script.extract()

        body_elm = soup.find("div", {"id": "mw-content-text"})
        title_elm = soup.find("span", {"class": "mw-page-title-main"})

        main_title = None if soup.title is None else soup.title.string
        webpage_text = ""

        if body_elm:
            if title_elm and len(title_elm) > 0:
                main_title = title_elm.string
            webpage_text = f"# {main_title}\n\n" + _CustomMarkdownify().convert_soup(body_elm)
        else:
            webpage_text = _CustomMarkdownify().convert_soup(soup)

        return DocumentConverterResult(title=main_title, text_content=str(webpage_text))


class YouTubeConverter(DocumentConverter):
    """Converter for YouTube video pages."""
    
    def convert(self, local_path: str, **kwargs: Any) -> Optional[DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None
            
        url = kwargs.get("url", "")
        if not url.startswith("https://www.youtube.com/watch?"):
            return None

        with open(local_path, "rt", encoding="utf-8") as fh:
            soup = BeautifulSoup(fh.read(), "html.parser")

        metadata = {"title": soup.title.string if soup.title else ""}
        for meta in soup(["meta"]):
            for attr in meta.attrs:
                if attr in ["itemprop", "property", "name"]:
                    metadata[meta[attr]] = meta.get("content", "")
                    break

        self._extract_yt_initial_data(soup, metadata)

        webpage_text = self._build_youtube_content(metadata, url)
        title = metadata.get("title", soup.title.string if soup.title else "YouTube Video")
        
        return DocumentConverterResult(title=title, text_content=webpage_text)

    def _extract_yt_initial_data(self, soup: BeautifulSoup, metadata: Dict[str, str]) -> None:
        try:
            for script in soup(["script"]):
                content = script.text
                if "ytInitialData" in content:
                    lines = re.split(r"\r?\n", content)
                    obj_start = lines[0].find("{")
                    obj_end = lines[0].rfind("}")
                    if obj_start >= 0 and obj_end >= 0:
                        data = json.loads(lines[0][obj_start:obj_end + 1])
                        attrdesc = self._find_key(data, "attributedDescriptionBodyText")
                        if attrdesc:
                            metadata["description"] = str(attrdesc["content"])
                    break
        except Exception:
            pass

    def _build_youtube_content(self, metadata: Dict[str, str], url: str) -> str:
        content = "# YouTube\n"
        title = self._get_metadata_value(metadata, ["title", "og:title", "name"])
        
        if title:
            content += f"\n## {title}\n"

        stats = self._build_video_stats(metadata)
        if stats:
            content += f"\n### Video Metadata\n{stats}\n"

        description = self._get_metadata_value(metadata, ["description", "og:description"])
        if description:
            content += f"\n### Description\n{description}\n"

        transcript = self._get_transcript(url)
        if transcript:
            content += f"\n### Transcript\n{transcript}\n"

        return content

    def _build_video_stats(self, metadata: Dict[str, str]) -> str:
        stats = []
        if views := self._get_metadata_value(metadata, ["interactionCount"]):
            stats.append(f"- **Views:** {views}")
        if keywords := self._get_metadata_value(metadata, ["keywords"]):
            stats.append(f"- **Keywords:** {keywords}")
        if runtime := self._get_metadata_value(metadata, ["duration"]):
            stats.append(f"- **Runtime:** {runtime}")
        return "\n".join(stats)

    def _get_metadata_value(self, metadata: Dict[str, str], keys: List[str]) -> Optional[str]:
        for key in keys:
            if key in metadata:
                return metadata[key]
        return None

    def _find_key(self, data: Any, target_key: str) -> Optional[Any]:
        if isinstance(data, list):
            for item in data:
                result = self._find_key(item, target_key)
                if result is not None:
                    return result
        elif isinstance(data, dict):
            for key, value in data.items():
                if key == target_key:
                    return value
                result = self._find_key(value, target_key)
                if result is not None:
                    return result
        return None

    def _get_transcript(self, url: str) -> str:
        try:
            parsed_url = urlparse(url)
            params = parse_qs(parsed_url.query)
            if "v" in params:
                video_id = str(params["v"][0])
                transcript = YouTubeTranscriptApi.get_transcript(video_id)
                return SRTFormatter().format_transcript(transcript)
        except Exception:
            pass
        return ""


class PdfConverter(DocumentConverter):
    """Converter for PDF files."""
    
    def convert(self, local_path: str, **kwargs: Any) -> Optional[DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pdf":
            return None

        text_content = pdfminer.high_level.extract_text(local_path)
        return DocumentConverterResult(title=None, text_content=text_content)


class DocxConverter(HtmlConverter):
    """Converter for DOCX files."""
    
    def convert(self, local_path: str, **kwargs: Any) -> Optional[DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".docx":
            return None

        with open(local_path, "rb") as docx_file:
            result = mammoth.convert_to_html(docx_file)
            return self._convert(result.value)


class XlsxConverter(HtmlConverter):
    """Converter for Excel files."""
    
    def convert(self, local_path: str, **kwargs: Any) -> Optional[DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".xlsx", ".xls"]:
            return None

        sheets = pd.read_excel(local_path, sheet_name=None)
        md_content = ""
        
        for sheet_name, sheet_data in sheets.items():
            md_content += f"## {sheet_name}\n"
            html_content = sheet_data.to_html(index=False)
            md_content += self._convert(html_content).text_content.strip() + "\n\n"

        return DocumentConverterResult(title=None, text_content=md_content.strip())


class PptxConverter(HtmlConverter):
    """Converter for PowerPoint files."""
    
    def convert(self, local_path: str, **kwargs: Any) -> Optional[DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pptx":
            return None

        presentation = pptx.Presentation(local_path)
        md_content = ""
        slide_num = 0

        for slide in presentation.slides:
            slide_num += 1
            md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"
            md_content += self._process_slide(slide)

        return DocumentConverterResult(title=None, text_content=md_content.strip())

    def _process_slide(self, slide: Any) -> str:
        content = ""
        title = slide.shapes.title

        for shape in slide.shapes:
            if self._is_picture(shape):
                content += self._process_picture(shape)
            elif self._is_table(shape):
                content += self._process_table(shape)
            elif shape.has_text_frame:
                if shape == title:
                    content += f"# {shape.text.lstrip()}\n"
                else:
                    content += shape.text + "\n"

        content = content.strip()

        if slide.has_notes_slide:
            content += "\n\n### Notes:\n"
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame is not None:
                content += notes_frame.text
            content = content.strip()

        return content

    def _is_picture(self, shape: Any) -> bool:
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if (shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER 
            and hasattr(shape, "image")):
            return True
        return False

    def _is_table(self, shape: Any) -> bool:
        return shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE

    def _process_picture(self, shape: Any) -> str:
        alt_text = ""
        try:
            alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
        except Exception:
            alt_text = shape.name

        filename = re.sub(r"\W", "", shape.name) + ".jpg"
        return f"\n![{alt_text}]({filename})\n"

    def _process_table(self, shape: Any) -> str:
        html_table = "<html><body><table>"
        first_row = True
        
        for row in shape.table.rows:
            html_table += "<tr>"
            for cell in row.cells:
                tag = "th" if first_row else "td"
                html_table += f"<{tag}>{html.escape(cell.text)}</{tag}>"
            html_table += "</tr>"
            first_row = False
            
        html_table += "</table></body></html>"
        return "\n" + self._convert(html_table).text_content.strip() + "\n"


class MediaConverter(DocumentConverter):
    """Base class for media file converters with metadata extraction."""
    
    def _get_metadata(self, local_path: str) -> Optional[Dict[str, Any]]:
        exiftool = shutil.which("exiftool")
        if not exiftool:
            return None

        try:
            result = subprocess.run([exiftool, "-json", local_path], 
                                  capture_output=True, text=True).stdout
            return json.loads(result)[0]
        except Exception:
            return None


class WavConverter(MediaConverter):
    """Converter for WAV audio files."""
    
    def convert(self, local_path: str, **kwargs: Any) -> Optional[DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".wav":
            return None

        md_content = self._extract_metadata(local_path)
        transcript = self._transcribe_audio(local_path)
        md_content += f"\n\n### Audio Transcript:\n{transcript}"

        return DocumentConverterResult(title=None, text_content=md_content.strip())

    def _extract_metadata(self, local_path: str) -> str:
        metadata = self._get_metadata(local_path)
        if not metadata:
            return ""

        content = ""
        fields = ["Title", "Artist", "Author", "Band", "Album", "Genre", 
                 "Track", "DateTimeOriginal", "CreateDate", "Duration"]
        
        for field in fields:
            if field in metadata:
                content += f"{field}: {metadata[field]}\n"
                
        return content

    def _transcribe_audio(self, local_path: str) -> str:
        try:
            recognizer = sr.Recognizer()
            with sr.AudioFile(local_path) as source:
                audio = recognizer.record(source)
                transcript = recognizer.recognize_google(audio).strip()
                return transcript if transcript else "[No speech detected]"
        except Exception:
            return "Error: Could not transcribe this audio."


class Mp3Converter(WavConverter):
    """Converter for MP3 and M4A audio files."""
    
    def convert(self, local_path: str, **kwargs: Any) -> Optional[DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".mp3", ".m4a"]:
            return None

        md_content = self._extract_metadata(local_path)
        transcript = self._transcribe_audio_file(local_path, extension)
        md_content += f"\n\n### Audio Transcript:\n{transcript}"

        return DocumentConverterResult(title=None, text_content=md_content.strip())

    def _transcribe_audio_file(self, local_path: str, extension: str) -> str:
        handle, temp_path = tempfile.mkstemp(suffix=".wav")
        os.close(handle)
        
        try:
            if extension.lower() == ".mp3":
                sound = pydub.AudioSegment.from_mp3(local_path)
            else:
                sound = pydub.AudioSegment.from_file(local_path, format="m4a")
            sound.export(temp_path, format="wav")
            
            return super()._transcribe_audio(temp_path)
        except Exception:
            return "Error: Could not transcribe this audio."
        finally:
            os.unlink(temp_path)


class ZipConverter(DocumentConverter):
    """Converter for ZIP archive files."""
    
    def __init__(self, extract_dir: str = "downloads"):
        self.extract_dir = extract_dir
        os.makedirs(self.extract_dir, exist_ok=True)

    def convert(self, local_path: str, **kwargs: Any) -> Optional[DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".zip" or not zipfile.is_zipfile(local_path):
            return None

        extracted_files = []
        with zipfile.ZipFile(local_path, "r") as zip_ref:
            zip_ref.extractall(self.extract_dir)
            for file_path in zip_ref.namelist():
                if not file_path.endswith("/"):
                    extracted_files.append(os.path.join(self.extract_dir, file_path))

        extracted_files.sort()
        md_content = "Downloaded the following files:\n" + "\n".join(f"* {f}" for f in extracted_files)

        return DocumentConverterResult(title="Extracted Files", text_content=md_content.strip())


class ImageConverter(MediaConverter):
    """Converter for image files with optional MLM description."""
    
    def convert(self, local_path: str, **kwargs: Any) -> Optional[DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".jpg", ".jpeg", ".png"]:
            return None

        md_content = self._extract_metadata(local_path)
        
        mlm_client = kwargs.get("mlm_client")
        mlm_model = kwargs.get("mlm_model")
        if mlm_client and mlm_model:
            description = self._get_mlm_description(local_path, extension, mlm_client, 
                                                  mlm_model, kwargs.get("mlm_prompt"))
            md_content += f"\n# Description:\n{description.strip()}\n"

        return DocumentConverterResult(title=None, text_content=md_content)

    def _extract_metadata(self, local_path: str) -> str:
        metadata = self._get_metadata(local_path)
        if not metadata:
            return ""

        content = ""
        fields = ["ImageSize", "Title", "Caption", "Description", "Keywords", 
                 "Artist", "Author", "DateTimeOriginal", "CreateDate", "GPSPosition"]
        
        for field in fields:
            if field in metadata:
                content += f"{field}: {metadata[field]}\n"
                
        return content

    def _get_mlm_description(self, local_path: str, extension: str, client: Any, 
                           model: str, prompt: Optional[str] = None) -> str:
        if not prompt or prompt.strip() == "":
            prompt = "Write a detailed caption for this image."

        with open(local_path, "rb") as image_file:
            content_type = mimetypes.guess_type("_dummy" + extension)[0] or "image/jpeg"
            image_base64 = base64.b64encode(image_file.read()).decode("utf-8")
            data_uri = f"data:{content_type};base64,{image_base64}"

        messages = [{
            "role": "user",
            "content": [
                {"type": "text", "text": prompt},
                {"type": "image_url", "image_url": {"url": data_uri}},
            ]
        }]

        response = client.chat.completions.create(model=model, messages=messages)
        return response.choices[0].message.content or ""


class FileConversionException(Exception):
    """Exception raised when file conversion fails."""
    pass


class UnsupportedFormatException(Exception):
    """Exception raised when file format is not supported."""
    pass


class MarkdownConverter:
    """Main converter class that handles multiple file formats."""
    
    def __init__(
        self,
        requests_session: Optional[requests.Session] = None,
        mlm_client: Optional[Any] = None,
        mlm_model: Optional[Any] = None,
    ):
        self._requests_session = requests_session or requests.Session()
        self._mlm_client = mlm_client
        self._mlm_model = mlm_model
        self._page_converters: List[DocumentConverter] = []
        
        self._register_default_converters()

    def _register_default_converters(self) -> None:
        converters = [
            PlainTextConverter(), HtmlConverter(), WikipediaConverter(),
            YouTubeConverter(), DocxConverter(), XlsxConverter(),
            PptxConverter(), WavConverter(), Mp3Converter(),
            ImageConverter(), ZipConverter(), PdfConverter()
        ]
        
        for converter in converters:
            self.register_page_converter(converter)

    def convert(self, source: Union[str, requests.Response], **kwargs: Any) -> DocumentConverterResult:
        if isinstance(source, str):
            if source.startswith(("http://", "https://", "file://")):
                return self.convert_url(source, **kwargs)
            else:
                return self.convert_local(source, **kwargs)
        elif isinstance(source, requests.Response):
            return self.convert_response(source, **kwargs)
        else:
            raise ValueError("Unsupported source type")

    def convert_local(self, path: str, **kwargs: Any) -> DocumentConverterResult:
        extensions = self._get_extensions(path, kwargs.get("file_extension"))
        return self._convert(path, extensions, **kwargs)

    def convert_stream(self, stream: Any, **kwargs: Any) -> DocumentConverterResult:
        extensions = [kwargs.get("file_extension")]
        
        handle, temp_path = tempfile.mkstemp()
        try:
            with os.fdopen(handle, "wb") as fh:
                content = stream.read()
                fh.write(content.encode("utf-8") if isinstance(content, str) else content)
            
            extensions.extend(self._get_extensions(temp_path, None))
            return self._convert(temp_path, extensions, **kwargs)
        finally:
            os.unlink(temp_path)

    def convert_url(self, url: str, **kwargs: Any) -> DocumentConverterResult:
        user_agent = "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
        timeout = (20, 40)
        
        response = self._requests_session.get(
            url, stream=True, 
            headers={"User-Agent": user_agent}, 
            timeout=timeout
        )
        response.raise_for_status()
        return self.convert_response(response, **kwargs)

    def convert_response(self, response: requests.Response, **kwargs: Any) -> DocumentConverterResult:
        extensions = self._get_extensions_from_response(response, kwargs.get("file_extension"))
        
        handle, temp_path = tempfile.mkstemp()
        try:
            with os.fdopen(handle, "wb") as fh:
                for chunk in response.iter_content(chunk_size=512):
                    fh.write(chunk)
            
            extensions.extend(self._get_extensions(temp_path, None))
            return self._convert(temp_path, extensions, url=response.url, **kwargs)
        except Exception as e:
            raise FileConversionException(f"Error in converting: {e}")
        finally:
            os.unlink(temp_path)

    def _convert(self, local_path: str, extensions: List[Optional[str]], **kwargs: Any) -> DocumentConverterResult:
        error_trace = ""
        
        for ext in extensions + [None]:
            for converter in self._page_converters:
                converter_kwargs = self._prepare_converter_kwargs(kwargs, ext)
                
                try:
                    result = converter.convert(local_path, **converter_kwargs)
                    if result is not None:
                        return self._normalize_result(result)
                except Exception:
                    error_trace = ("\n\n" + traceback.format_exc()).strip()

        if error_trace:
            raise FileConversionException(
                f"Could not convert '{local_path}' to Markdown. "
                f"File type was recognized as {extensions}. Error:\n{error_trace}"
            )
            
        raise UnsupportedFormatException(
            f"Could not convert '{local_path}' to Markdown. "
            f"The formats {extensions} are not supported."
        )

    def _prepare_converter_kwargs(self, base_kwargs: Dict[str, Any], extension: Optional[str]) -> Dict[str, Any]:
        kwargs = copy.deepcopy(base_kwargs)
        
        if extension is not None:
            kwargs["file_extension"] = extension
        elif "file_extension" in kwargs:
            del kwargs["file_extension"]

        if self._mlm_client and "mlm_client" not in kwargs:
            kwargs["mlm_client"] = self._mlm_client
        if self._mlm_model and "mlm_model" not in kwargs:
            kwargs["mlm_model"] = self._mlm_model
            
        return kwargs

    def _normalize_result(self, result: DocumentConverterResult) -> DocumentConverterResult:
        result.text_content = "\n".join(line.rstrip() for line in re.split(r"\r?\n", result.text_content))
        result.text_content = re.sub(r"\n{3,}", "\n\n", result.text_content)
        return result

    def _get_extensions(self, path: str, initial_ext: Optional[str]) -> List[Optional[str]]:
        extensions = [initial_ext] if initial_ext is not None else []
        
        base, ext = os.path.splitext(path)
        self._append_unique_extension(extensions, ext)
        self._append_unique_extension(extensions, self._guess_extension_magic(path))
        
        return extensions

    def _get_extensions_from_response(self, response: requests.Response, initial_ext: Optional[str]) -> List[Optional[str]]:
        extensions = [initial_ext] if initial_ext is not None else []
        
        content_type = response.headers.get("content-type", "").split(";")[0]
        self._append_unique_extension(extensions, mimetypes.guess_extension(content_type))

        content_disposition = response.headers.get("content-disposition", "")
        if match := re.search(r"filename=([^;]+)", content_disposition):
            base, ext = os.path.splitext(match.group(1).strip("\"'"))
            self._append_unique_extension(extensions, ext)

        base, ext = os.path.splitext(urlparse(response.url).path)
        self._append_unique_extension(extensions, ext)
        
        return extensions

    def _append_unique_extension(self, extensions: List[Optional[str]], ext: Optional[str]) -> None:
        if ext and ext.strip() and ext not in extensions:
            extensions.append(ext.strip())

    def _guess_extension_magic(self, path: str) -> Optional[str]:
        try:
            guesses = puremagic.magic_file(path)
            if guesses and (ext := guesses[0].extension.strip()):
                return ext
        except (FileNotFoundError, IsADirectoryError, PermissionError):
            pass
        return None

    def register_page_converter(self, converter: DocumentConverter) -> None:
        """Register a document converter (higher priority converters should be registered last)."""
        self._page_converters.insert(0, converter)