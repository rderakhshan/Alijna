#!/usr/bin/env python
# coding=utf-8
# Copyright 2025 The OPPO Inc. Personal AI team. All rights reserved.
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

import base64
import json
import mimetypes
import os
import uuid
from typing import Optional
import openai

import requests
from dotenv import load_dotenv
from PIL import Image

from .tools import Tool
from .models import Model, MessageRole
from .mm_tools_utils import MarkdownConverter

from xml.dom import minidom
from openpyxl import load_workbook
import json
from Bio import PDB
from pptx import Presentation
import shutil
import textwrap

load_dotenv(override=True)

class VisualInspectorTool(Tool):
    name = "inspect_file_as_image"
    description = """
You cannot load files directly: use this tool to process image files and answer related questions.
This tool supports the following image formats: [".jpg", ".jpeg", ".png", ".gif", ".bmp"]. For other file types, use the appropriate inspection tool."""

    inputs = {
        "file_path": {
            "description": "The path to the file you want to read as an image. Must be a '.something' file, like '.jpg','.png','.gif'. If it is text, use the text_inspector tool instead! If it is audio, use the audio_inspector tool instead! DO NOT use this tool for an HTML webpage: use the web_search tool instead!",
            "type": "string",
        },
        "question": {
            "description": "[Optional]: Your question about the image content. Provide as much context as possible. Do not pass this parameter if you just want to get a description of the image.",
            "type": "string",
            "nullable": True,
        },
    }
    output_type = "string"

    def __init__(self, model: Model, text_limit: int):
        super().__init__()
        self.model = model
        self.text_limit = text_limit
        self.gpt_key = os.getenv("OPENAI_API_KEY")
        self.gpt_url = os.getenv("OPENAI_BASE_URL")

    def _validate_file_type(self, file_path: str):
        if not any(file_path.lower().endswith(ext) for ext in [".jpg", ".jpeg", ".png", ".gif", ".bmp"]):
            raise ValueError("Unsupported file type. Use the appropriate tool for text/audio files.")

    def _resize_image(self, image_path: str) -> str:
        img = Image.open(image_path)
        width, height = img.size
        img = img.resize((int(width / 2), int(height / 2)))
        new_image_path = f"resized_{os.path.basename(image_path)}"
        img.save(new_image_path)
        return new_image_path

    def _encode_image(self, image_path: str) -> str:
        if image_path.startswith("http"):
            user_agent = "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/119.0.0.0 Safari/537.36 Edg/119.0.0.0"
            request_kwargs = {
                "headers": {"User-Agent": user_agent},
                "stream": True,
            }

            response = requests.get(image_path, **request_kwargs)
            response.raise_for_status()
            content_type = response.headers.get("content-type", "")

            extension = mimetypes.guess_extension(content_type)
            if extension is None:
                extension = ".download"

            fname = str(uuid.uuid4()) + extension
            download_path = os.path.abspath(os.path.join("downloads", fname))

            with open(download_path, "wb") as fh:
                for chunk in response.iter_content(chunk_size=512):
                    fh.write(chunk)

            image_path = download_path

        with open(image_path, "rb") as image_file:
            return base64.b64encode(image_file.read()).decode("utf-8")


    def forward(self, file_path: str, question: Optional[str] = None) -> str:
        self._validate_file_type(file_path)
        
        if not question:
            question = "Please write a detailed caption for this image."
        try:
            mime_type, _ = mimetypes.guess_type(file_path)
            base64_image = self._encode_image(file_path)
            payload = {
                "model": "gpt-4o",
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": question},
                            {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{base64_image}"}},
                        ],
                    }
                ],
                "max_tokens": 2000,
                "top_p": 0.1,
            }

            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {self.gpt_key}"
            }

            response = requests.post(
                f"{self.gpt_url}/chat/completions",
                headers=headers,
                json=payload
            )
            response.raise_for_status()
            description = response.json()["choices"][0]["message"]["content"]
        except Exception as gpt_error:
            return f"Visual processing failed: {str(gpt_error)}"

        if not question.startswith("Please write a detailed caption"):
            return description
        return f"You did not provide a particular question, so here is a detailed description of the image: {description}"
    

class TextInspectorTool(Tool):
    name = "inspect_file_as_text"
    description = """
You cannot load files yourself: instead call this tool to read a file as markdown text and ask questions about it.
This tool handles the following file extensions: [".html", ".pdb", ".xlsx", ".xls", ".pdf", ".docx", ".ppt", ".pptx"], and all other types of text files. IT DOES NOT HANDLE IMAGES."""

    inputs = {
        "file_path": {
            "description": "The path to the file you want to read as text. Must be a '.something' file, like '.pdf'. If it is an image, use the visualizer tool instead! If it is an audio, use the audio tool instead! DO NOT use this tool for an HTML webpage: use the web_search tool instead!",
            "type": "string",
        },
        "question": {
            "description": "[Optional]: Your question, as a natural language sentence. Provide as much context as possible. Do not pass this parameter if you just want to directly return the content of the file.",
            "type": "string",
            "nullable": True,
        },
    }
    output_type = "string"
    md_converter = MarkdownConverter()

    def __init__(self, model: Model, text_limit: int):
        super().__init__()
        self.model = model
        self.text_limit = text_limit

    def jsonld_to_markdown(self, data):
        markdown = ""
        if isinstance(data, dict):
            for key, value in data.items():
                markdown += f"**{key}**: {self.jsonld_to_markdown(value)}\n"
        elif isinstance(data, list):
            for item in data:
                markdown += f"- {self.jsonld_to_markdown(item)}\n"
        else:
            markdown += str(data)
        return markdown
    
    def parse_pdb_file(self, file_path):

        parser = PDB.PDBParser(QUIET=True)
        structure = parser.get_structure("protein", file_path)

        atoms = list(structure.get_atoms())
        if len(atoms) < 2:
            return "Error: PDB file contains fewer than two atoms."

        atom1, atom2 = atoms[0], atoms[1]
        distance = atom1 - atom2

        return f"First atom: {atom1.get_name()} ({atom1.coord})\n" \
            f"Second atom: {atom2.get_name()} ({atom2.coord})\n" \
            f"Distance: {distance:.3f} Angstroms ({distance * 100:.0f} pm)"

    def extract_excel_data(self, file_path, max_rows=30):
        try:
            workbook = load_workbook(file_path)
            result = []

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]

                for row in sheet.iter_rows():
                    row_data = []
                    for cell in row:
                        cell_value = cell.value if cell.value is not None else ""

                        fill = cell.fill
                        if hasattr(fill, "fgColor") and fill.fgColor.type == "rgb":
                            cell_color = fill.fgColor.rgb
                            if len(cell_color) == 8:
                                cell_color = cell_color[2:]
                        else:
                            cell_color = "FFFFFF"
                        row_data.append({
                            "value": cell_value,
                            "color": cell_color
                        })
                    result.append(row_data)
                json_data = json.dumps(result, ensure_ascii=False, indent=4)
                data_with_colors = json.loads(json_data)
                text = []
                num_rows = len(data_with_colors)
                num_cols = len(data_with_colors[0]) if data_with_colors else 0
                text.append(f"This is a {num_rows} rows and {num_cols} columns table. The content is shown below:")
                for i, row in enumerate(data_with_colors):
                    row_text = ""
                    for j, cell in enumerate(row):
                        if cell["value"] != "":
                            value = cell["value"]
                        else:
                            value = "None"
                        color = cell["color"]
                        if color == "FFFFFF" or color == "000000":
                            row_text += f"{value} "
                        else:
                            row_text += f"{value}({color}) "
                    text.append(row_text)
            return "\n".join(text)
        except Exception as e:
            return f"Error: {str(e)}"

    def forward_initial_exam_mode(self, file_path, question):
        result = self.md_converter.convert(file_path)

        if file_path.endswith(".xml"):
            try:
                dom = minidom.parse(file_path)
                result_text = " ".join(
                    [node.firstChild.nodeValue for node in dom.getElementsByTagName("*") if node.firstChild and node.firstChild.nodeType == node.TEXT_NODE]
                )
                result.text_content = result_text
            except Exception as e:
                raise Exception(f"Error parsing XML file: {str(e)}")

        if file_path.endswith(".csv"):
            try:
                with open(file_path, 'r') as fr:
                    contents = fr.readlines()
                result.text_content = contents
            except Exception as e:
                raise Exception(f"Error parsing CSV file: {str(e)}")

        if file_path.endswith(".pdb"):
            try:
                pdb_info = self.parse_pdb_file(file_path)
                if not question:
                    return f"Extracted PDB Data:\n{pdb_info}"

                else:
                    result.text_content = pdb_info
            except Exception as e:
                raise Exception(f"Error parsing PDB file: {e}")

        if file_path.endswith((".ppt",  ".pptx")):
            content = ""
            try:
                ppt = Presentation(file_path)
                for slide_number, slide in enumerate(ppt.slides, start=1):
                    content += f"=== Slide {slide_number} ===\n"
                    slide_texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_texts.append(shape.text.strip())
                    content += '\n'.join(slide_texts) + '\n\n'
                result.text_content = content.strip()
            except Exception as e:
                return f"Error parsing PPT file: {e}"

        if file_path.endswith((".xls", ".xlsx")):
            try:
                full_content = self.extract_excel_data(file_path)
                if not question:
                    return full_content
                else:
                    result.text_content = full_content
            except Exception as e:
                raise Exception(f"Error processing Excel file: {e}")

        if file_path[-4:] in [".png", ".jpg"]:
            raise Exception("Cannot use inspect_file_as_text tool with images: use visualizer instead!")

        if ".zip" in file_path:
            return result.text_content
            
        if not question:
            return result.text_content

        if len(result.text_content) < 4000:
            return "Document content: " + result.text_content

        messages = [
            {
                "role": MessageRole.SYSTEM,
                "content": [
                    {
                        "type": "text",
                        "text": "Here is a file:\n### "
                        + str(result.title)
                        + "\n\n"
                        + result.text_content[: self.text_limit],
                    }
                ],
            },
            {
                "role": MessageRole.USER,
                "content": [
                    {
                        "type": "text",
                        "text": "Now please write a short, 5 sentence caption for this document, that could help someone asking this question: "
                        + question
                        + "\n\nDon't answer the question yourself! Just provide useful notes on the document",
                    }
                ],
            },
        ]
        return self.model(messages).content

    def forward(self, file_path, question: Optional[str] = None) -> str:

        if ".jsonld" in file_path:
            try:
                with open(file_path, 'r') as f:
                    data = json.load(f)
                result_text = self.jsonld_to_markdown(data)
                result = type('Result', (object,), {'title': file_path, 'text_content': result_text})
            except Exception as e:
                raise Exception(f"Error parsing JSON-LD file: {str(e)}")
        else:
            result = self.md_converter.convert(file_path)

        if file_path[-4:] in [".png", ".jpg"]:
            raise Exception("Cannot use inspect_file_as_text tool with images: use visualizer instead!")

        if ".zip" in file_path:
            return result.text_content

        if file_path.endswith((".ppt",  ".pptx")):
            content = ""
            try:
                ppt = Presentation(file_path)
                for slide_number, slide in enumerate(ppt.slides, start=1):
                    content += f"=== Slide {slide_number} ===\n"
                    slide_texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_texts.append(shape.text.strip())
                    content += '\n'.join(slide_texts) + '\n\n'
                result.text_content = content.strip()
            except Exception as e:
                return f"Error parsing PPT file: {e}"

        if file_path.endswith((".xls", ".xlsx")):
            try:
                full_content = self.extract_excel_data(file_path)
                if not question:
                    return full_content
                else:
                    result.text_content = full_content
            except Exception as e:
                raise Exception(f"Error processing Excel file: {e}")

        if file_path.endswith(".pdb"):
            try:
                pdb_info = self.parse_pdb_file(file_path)
                if not question:
                    return f"Extracted PDB Data:\n{pdb_info}"

                else:
                    result.text_content = pdb_info
            except Exception as e:
                raise Exception(f"Error parsing PDB file: {e}")
            
        if not question:
            return result.text_content

        messages = [
            {
                "role": MessageRole.SYSTEM,
                "content": [
                    {
                        "type": "text",
                        "text": "You will have to write a short caption for this file, then answer this question:"
                        + question,
                    }
                ],
            },
            {
                "role": MessageRole.USER,
                "content": [
                    {
                        "type": "text",
                        "text": "Here is the complete file:\n### "
                        + str(result.title)
                        + "\n\n"
                        + result.text_content[: self.text_limit],
                    }
                ],
            },
            {
                "role": MessageRole.USER,
                "content": [
                    {
                        "type": "text",
                        "text": "Now answer the question below. Use these three headings: '1. Short answer', '2. Extremely detailed answer', '3. Additional Context on the document and question asked'."
                        + question,
                    }
                ],
            },
        ]
        return self.model(messages).content
    
class AudioInspectorTool(Tool):
    name = "inspect_file_as_audio"
    description = """
You cannot load files directly: use this tool to process audio files and answer related questions.
This tool supports the following audio formats: [".mp3", ".m4a", ".wav"]. For other file types, use the appropriate inspection tool."""

    inputs = {
        "file_path": {
            "description": "The path to the file you want to read as audio. Must be a '.something' file, like '.mp3','.m4a','.wav'. If it is an text, use the text_inspector tool instead! If it is an image, use the visualizer tool instead! DO NOT use this tool for an HTML webpage: use the web_search tool instead!",
            "type": "string",
        },
        "question": {
            "description": "[Optional]: Your question about the audio content. Provide as much context as possible. Do not pass this parameter if you just want to directly return the content of the file.",
            "type": "string",
            "nullable": True,
        },
    }
    output_type = "string"

    def __init__(self, model: Model, text_limit: int):
        super().__init__()
        self.model = model
        self.text_limit = text_limit
        self.api_key = os.getenv("MTU_API_KEY")
        self.base_url = os.getenv("MTU_BASE_URL")

    def _validate_file_type(self, file_path: str):
        """Validate if the file type is a supported audio format"""
        if not any(file_path.endswith(ext) for ext in [".mp3", ".m4a", ".wav"]):
            raise ValueError("Unsupported file type. Use the appropriate tool for text/image files.")

    def transcribe_audio(self, file_path: str) -> str:
        """Transcribe audio using OpenAI Whisper API"""
        client = openai.OpenAI(api_key=self.api_key, base_url=self.base_url)
        try:
            with open(file_path, "rb") as audio_file:
                transcription = client.audio.transcriptions.create(
                    model="whisper-1",
                    file=audio_file
                )
            return transcription.text
        except Exception as e:
            raise RuntimeError(f"Speech recognition failed: {str(e)}") from e

    def forward(self, file_path: str, question: Optional[str] = None) -> str:
        self._validate_file_type(file_path)
        
        try:
            transcript = self.transcribe_audio(file_path)
        except Exception as e:
            return f"Audio processing error: {str(e)}"
        
        if not question:
            return f"Audio transcription:\n{transcript[:self.text_limit]}"

        messages = [
            {
                "role": MessageRole.SYSTEM,
                "content": [{
                    "type": "text",
                    "text": f"Audio transcription:\n{transcript[:self.text_limit]}"
                }]
            },
            {
                "role": MessageRole.USER,
                "content": [{
                    "type": "text",
                    "text": f"Answer the following question based on the audio content using the format:\n"
                            "1. Brief answer\n2. Detailed analysis\n3. Relevant context\nQuestion: {question}"
                }]
            }
        ]
        
        return self.model(messages).content

def get_image_description(file_name: str, question: str, visual_inspection_tool) -> str:
    prompt = f"""Write a caption of 5 sentences for this image. Pay special attention to any details that might be useful for someone answering the following question:
{question}. But do not try to answer the question directly!
Do not add any information that is not present in the image."""
    return visual_inspection_tool(file_path=file_name, question=prompt)

def get_document_description(file_path: str, question: str, document_inspection_tool) -> str:
    prompt = f"""Write a caption of 5 sentences for this document. Pay special attention to any details that might be useful for someone answering the following question:
{question}. But do not try to answer the question directly!
Do not add any information that is not present in the document."""
    return document_inspection_tool.forward_initial_exam_mode(file_path=file_path, question=prompt)

def get_audio_description(audio_path: str, question: str, audio_inspection_tool) -> str:
    prompt = f"""Write a caption of 5 sentences for this audio. Pay special attention to any details that might be useful for someone answering the following question:
{question}. But do not try to answer the question directly!
Do not add any information that is not present in the audio."""
    return audio_inspection_tool.forward(file_path=audio_path, question=prompt)


def get_single_file_description(file_path: str, question: str, visual_inspection_tool, document_inspection_tool, audio_inspection_tool):
    file_extension = file_path.split(".")[-1]
    if file_extension in ["png", "jpg", "jpeg"]:
        file_description = f" - Attached image: {file_path}"
        file_description += (
            f"\n     -> Image description: {get_image_description(file_path, question, visual_inspection_tool)}"
        )
        return file_description
    elif file_extension in ["pdf", "xls", "xlsx", "docx", "doc", "xml", "ppt"]:
        file_description = f" - Attached document: {file_path}"
        image_path = file_path.split(".")[0] + ".png"
        description = get_document_description(file_path, question, document_inspection_tool)
        file_description += f"\n     -> File description: {description}"
        if os.path.exists(image_path):
            description = get_image_description(image_path, question, visual_inspection_tool)
            file_description += f"\n     -> Additional image description: {description}"
        return file_description
    elif file_extension in ["mp3", "m4a", "wav"]:
        file_description = f" - Attached audio: {file_path}"
        file_description += (
            f"\n     -> File description: {get_audio_description(file_path, question, audio_inspection_tool)}"
        )
        return file_description
    else:
        return f" - Attached file: {file_path}"


def get_zip_description(file_path: str, question: str, visual_inspection_tool, document_inspection_tool, audio_inspection_tool):
    folder_path = file_path.replace(".zip", "")
    os.makedirs(folder_path, exist_ok=True)
    shutil.unpack_archive(file_path, folder_path)

    prompt_use_files = ""
    for root, dirs, files in os.walk(folder_path):
        for file in files:
            file_path = os.path.join(root, file)
            prompt_use_files += "\n" + textwrap.indent(
                get_single_file_description(file_path, question, visual_inspection_tool, document_inspection_tool, audio_inspection_tool),
                prefix="    ",
            )
    return prompt_use_files